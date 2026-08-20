#!/usr/bin/env python3
"""
Merge multiple single-slide PPTX files into one presentation.

Usage:
    python scripts/merge_slides.py slide_01.pptx slide_02.pptx ... -o final.pptx

This is a backup mechanism for when the single-JS-file approach hits
token or execution limits. Normally, all slides are generated in one
JS file and this script is not needed.

Requires: python-pptx
    pip install python-pptx
"""

import argparse
import copy
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def merge_presentations(input_files: list[str], output_file: str) -> None:
    """Merge multiple PPTX files into a single presentation.

    Each input file contributes its slides to the output in order.
    The first file's slide dimensions are used for the output.
    """
    if not input_files:
        print("Error: No input files provided")
        sys.exit(1)

    # Validate all input files exist
    for f in input_files:
        if not os.path.exists(f):
            print(f"Error: File not found: {f}")
            sys.exit(1)

    # Use the first presentation as the base
    base_prs = Presentation(input_files[0])
    slide_width = base_prs.slide_width
    slide_height = base_prs.slide_height

    print(f"Base dimensions: {slide_width}x{slide_height} EMU")
    print(f"  = {slide_width / 914400:.1f}\" x {slide_height / 914400:.1f}\"")
    print(f"Merging {len(input_files)} files...")

    # For the first file, slides are already in base_prs
    slide_count = len(base_prs.slides)
    print(f"  [1/{len(input_files)}] {input_files[0]}: {slide_count} slide(s)")

    # Merge subsequent files
    for idx, input_file in enumerate(input_files[1:], start=2):
        src_prs = Presentation(input_file)
        src_slide_count = len(src_prs.slides)

        for src_slide in src_prs.slides:
            # Add a blank slide using the first layout
            slide_layout = base_prs.slide_layouts[6]  # Blank layout
            new_slide = base_prs.slides.add_slide(slide_layout)

            # Copy all shapes from source to new slide
            for shape in src_slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            # Copy background if set
            if src_slide.background.fill.type is not None:
                bg_elem = copy.deepcopy(src_slide.background._element)
                new_slide.background._element.getparent().replace(
                    new_slide.background._element, bg_elem
                )

            slide_count += 1

        print(f"  [{idx}/{len(input_files)}] {input_file}: {src_slide_count} slide(s)")

    # Save merged presentation
    base_prs.save(output_file)
    print(f"\nMerged {slide_count} slides into: {output_file}")


def main():
    parser = argparse.ArgumentParser(
        description="Merge multiple PPTX files into one presentation"
    )
    parser.add_argument("inputs", nargs="+", help="Input PPTX files to merge (in order)")
    parser.add_argument(
        "-o", "--output", default="merged.pptx", help="Output file path (default: merged.pptx)"
    )
    args = parser.parse_args()
    merge_presentations(args.inputs, args.output)


if __name__ == "__main__":
    main()
